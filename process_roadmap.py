import os
import json
import traceback
import requests
from datetime import datetime, timedelta
from docx import Document
from pptx import Presentation
from pptx.util import Inches
from openpyxl import load_workbook
from google.oauth2 import service_account
from googleapiclient.discovery import build
from googleapiclient.http import MediaFileUpload

# === Google Drive Setup ===
drive_service = None
try:
    creds_json = os.getenv("GOOGLE_SERVICE_ACCOUNT_JSON")
    if creds_json:
        creds = service_account.Credentials.from_service_account_info(
            json.loads(creds_json),
            scopes=["https://www.googleapis.com/auth/drive"]
        )
        drive_service = build("drive", "v3", credentials=creds)
except Exception as e:
    print(f"❌ Google Drive setup failed: {e}")
    traceback.print_exc()

def upload_to_drive(file_path, session_id):
    try:
        query = f"name='{session_id}' and mimeType='application/vnd.google-apps.folder'"
        results = drive_service.files().list(q=query, fields="files(id)").execute()
        folders = results.get("files", [])
        if folders:
            folder_id = folders[0]["id"]
        else:
            folder = drive_service.files().create(body={
                "name": session_id,
                "mimeType": "application/vnd.google-apps.folder"
            }, fields="id").execute()
            folder_id = folder["id"]

        file_meta = {"name": os.path.basename(file_path), "parents": [folder_id]}
        media = MediaFileUpload(file_path, resumable=True)
        uploaded = drive_service.files().create(body=file_meta, media_body=media, fields="id").execute()
        return f"https://drive.google.com/file/d/{uploaded['id']}/view"
    except Exception as e:
        print(f"❌ Upload failed: {e}")
        return None

def download_files(files, folder_path):
    downloaded = []
    for f in files:
        if not f.get("file_url"):
            print(f"⚠️ Skipping {f.get('file_name')} – missing URL")
            continue
        try:
            path = os.path.join(folder_path, f["file_name"])
            r = requests.get(f["file_url"], timeout=15)
            with open(path, "wb") as out:
                out.write(r.content)
            f["local_path"] = path
            downloaded.append(f)
        except Exception as e:
            print(f"❌ Download failed for {f['file_name']}: {e}")
    return downloaded

def extract_devices(files):
    devices = []
    for f in files:
        if f["file_type"] in ["gap_hw", "gap_sw"]:
            wb = load_workbook(f["local_path"])
            sheet = wb.active
            for row in sheet.iter_rows(min_row=2, values_only=True):
                device = {
                    "name": row[0],
                    "platform": row[2],
                    "tier": row[3],
                    "status": row[4],
                    "recommendation": row[5],
                }
                devices.append(device)
    return devices

def generate_roadmap_docx(session_id, folder_path, devices):
    path = os.path.join(folder_path, "IT_Transformation_Roadmap.docx")
    doc = Document()
    doc.add_heading("IT Transformation Roadmap", 0)
    doc.add_paragraph(f"Session: {session_id}").bold = True

    doc.add_heading("1. Executive Summary", level=1)
    doc.add_paragraph("This roadmap defines phases, timelines, epics, and configuration changes aligned with the organization's modernization goals.")

    doc.add_heading("2. Project Plan", level=1)
    doc.add_paragraph("Phased approach with estimated timelines:")

    doc.add_heading("3. Device Transformation Matrix", level=1)
    table = doc.add_table(rows=1, cols=5)
    hdr = table.rows[0].cells
    hdr[0].text = "Device"
    hdr[1].text = "Current Platform"
    hdr[2].text = "Tier"
    hdr[3].text = "Recommendation"
    hdr[4].text = "Status"
    for d in devices:
        row = table.add_row().cells
        row[0].text = d["name"]
        row[1].text = d["platform"]
        row[2].text = str(d["tier"])
        row[3].text = d["recommendation"] or "N/A"
        row[4].text = d["status"]

    doc.add_heading("4. Change Tickets", level=1)
    for i, d in enumerate(devices[:10]):  # limit for brevity
        doc.add_paragraph(f"""
        Change Ticket #{i+1}
        - CI: {d['name']}
        - Type: {"Normal" if "obsolete" in str(d['status']).lower() else "Standard"}
        - Justification: {d['recommendation'] or "Recommended upgrade"}
        - Schedule: Phase 1 (Month 1-2)
        - Impact: Moderate
        """)

    doc.add_heading("5. Agile Epics and Stories", level=1)
    doc.add_paragraph("Epic: Hardware Modernization")
    for i, d in enumerate(devices[:5]):
        doc.add_paragraph(f" - Story: Upgrade {d['name']} to meet performance targets.")

    doc.add_heading("6. Risk and Mitigation", level=1)
    doc.add_paragraph("Risk: Data loss during migration\nMitigation: Implement snapshot and rollback plan before all major hardware changes.")

    doc.save(path)
    return path

def generate_roadmap_pptx(session_id, folder_path, devices):
    path = os.path.join(folder_path, "IT_Transformation_Timeline.pptx")
    ppt = Presentation()
    slide = ppt.slides.add_slide(ppt.slide_layouts[0])
    slide.shapes.title.text = "Transformation Timeline"
    slide.placeholders[1].text = f"Session: {session_id}"

    def add_slide(title, bullets):
        s = ppt.slides.add_slide(ppt.slide_layouts[1])
        s.shapes.title.text = title
        tf = s.placeholders[1].text_frame
        tf.clear()
        for b in bullets:
            tf.add_paragraph().text = b

    add_slide("Phase Timeline", ["Phase 1: Infra upgrade (Month 1-3)", "Phase 2: App Modernization (Month 4-6)", "Phase 3: Cloud Migration (Month 7-12)"])
    add_slide("Epic Summary", [f"Upgrade {d['name']}" for d in devices[:5]])
    add_slide("CI Transitions", [f"{d['name']} → {d['recommendation'] or 'Optimized'}" for d in devices[:5]])

    ppt.save(path)
    return path

def process_roadmap(session_id, email, files, folder_path):
    try:
        os.makedirs(folder_path, exist_ok=True)
        downloaded = download_files(files, folder_path)
        devices = extract_devices(downloaded)

        docx_path = generate_roadmap_docx(session_id, folder_path, devices)
        pptx_path = generate_roadmap_pptx(session_id, folder_path, devices)

        docx_url = upload_to_drive(docx_path, session_id)
        pptx_url = upload_to_drive(pptx_path, session_id)

        for f in downloaded:
            f["file_url"] = upload_to_drive(f["local_path"], session_id)

        downloaded.extend([
            {
                "file_name": os.path.basename(docx_path),
                "file_url": docx_url,
                "file_type": "docx_roadmap"
            },
            {
                "file_name": os.path.basename(pptx_path),
                "file_url": pptx_url,
                "file_type": "pptx_roadmap"
            }
        ])

        NEXT_GPT_URL = "https://it-financials-api.onrender.com/start_it_financials"
        payload = {
            "session_id": session_id,
            "email": email,
            "gpt_module": "roadmap",
            "files": downloaded,
            "status": "complete"
        }
        requests.post(NEXT_GPT_URL, json=payload)

    except Exception as e:
        print(f"🔥 Roadmap processing failed: {e}")
        traceback.print_exc()
